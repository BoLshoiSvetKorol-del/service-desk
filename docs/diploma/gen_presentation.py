"""
Генератор презентации для дипломного проекта "Service Desk System".
Запуск: python gen_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

DIAGRAMS = os.path.join(os.path.dirname(__file__), "diagrams")

# ── Цветовая схема ──────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x1A, 0x23, 0x3A)   # тёмно-синий (фон заголовков)
C_ACCENT = RGBColor(0x21, 0x96, 0xF3)   # синий accent
C_LIGHT  = RGBColor(0xF5, 0xF7, 0xFA)   # светло-серый фон
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT   = RGBColor(0x1A, 0x23, 0x3A)
C_SUB    = RGBColor(0x55, 0x6A, 0x80)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # полностью пустой макет


# ── Вспомогательные функции ─────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=C_TEXT,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def slide_header(slide, title, subtitle=None):
    """Синяя полоса заголовка слева (вертикальная) + текст."""
    # Верхняя полоса
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), fill_color=C_DARK)
    add_text(slide, title,
             Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
             font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.4),
                 font_size=14, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.LEFT)
    # Тонкая линия-акцент под полосой
    add_rect(slide, 0, Inches(1.2), SLIDE_W, Pt(3), fill_color=C_ACCENT)


def add_image(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        # Заглушка если файл не найден
        r = add_rect(slide, left, top, width, height,
                     fill_color=RGBColor(0xDD, 0xDD, 0xDD))
        add_text(slide, f"[изображение:\n{os.path.basename(path)}]",
                 left + Inches(0.1), top + Inches(0.2),
                 width - Inches(0.2), height - Inches(0.4),
                 font_size=11, color=C_SUB, align=PP_ALIGN.CENTER)


def bullet_block(slide, items, left, top, width, height,
                 font_size=16, color=C_TEXT, bullet="•"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


# ════════════════════════════════════════════════════════════════════════════
# СЛАЙДЫ
# ════════════════════════════════════════════════════════════════════════════

# ── 1. ТИТУЛЬНЫЙ СЛАЙД ──────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
# Фон
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)
# Акцентная полоса слева
add_rect(s, 0, 0, Inches(0.15), SLIDE_H, fill_color=C_ACCENT)
# Декоративный прямоугольник
add_rect(s, Inches(0.4), Inches(4.7), Inches(12.5), Pt(2), fill_color=C_ACCENT)

add_text(s, "Дипломный проект",
         Inches(0.4), Inches(0.9), Inches(12), Inches(0.6),
         font_size=16, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.LEFT)

add_text(s, "Корпоративная система\nуправления заявками",
         Inches(0.4), Inches(1.5), Inches(12), Inches(2.2),
         font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

add_text(s, "Service Desk System",
         Inches(0.4), Inches(3.7), Inches(8), Inches(0.8),
         font_size=28, color=C_ACCENT, align=PP_ALIGN.LEFT)

add_text(s, "Специальность: Информационные системы и программирование\n"
            "Студент: Дегтярев М.Д.\n"
            "Руководитель: ___________________\n"
            "2026 г.",
         Inches(0.4), Inches(5.0), Inches(8), Inches(2.0),
         font_size=15, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.LEFT)


# ── 2. ЦЕЛИ И ЗАДАЧИ ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_LIGHT)
slide_header(s, "Цели и задачи дипломного проекта")

# Левая колонка — Цель
add_rect(s, Inches(0.3), Inches(1.45), Inches(5.9), Inches(5.7),
         fill_color=C_WHITE, line_color=RGBColor(0xDD, 0xEE, 0xFF), line_width=Pt(1))
add_rect(s, Inches(0.3), Inches(1.45), Inches(5.9), Inches(0.5), fill_color=C_ACCENT)
add_text(s, "ЦЕЛЬ", Inches(0.4), Inches(1.48), Inches(5.7), Inches(0.44),
         font_size=14, bold=True, color=C_WHITE)
add_text(s,
         "Разработка корпоративной веб-системы автоматизации "
         "процессов технической поддержки и управления заявками "
         "с поддержкой SLA, ролевого доступа и уведомлений в реальном времени.",
         Inches(0.5), Inches(2.05), Inches(5.5), Inches(1.5),
         font_size=15, color=C_TEXT)

tasks = [
    "Проанализировать предметную область управления ИТ-заявками",
    "Спроектировать архитектуру многоуровневого веб-приложения",
    "Разработать REST API на FastAPI с ролевой моделью доступа",
    "Реализовать SLA-движок с расчётом в рабочих часах",
    "Разработать React-интерфейс с real-time уведомлениями (SSE)",
    "Реализовать систему уведомлений (email + push) через Celery",
    "Провести тестирование и задокументировать систему",
]

add_rect(s, Inches(6.5), Inches(1.45), Inches(6.5), Inches(5.7),
         fill_color=C_WHITE, line_color=RGBColor(0xDD, 0xEE, 0xFF), line_width=Pt(1))
add_rect(s, Inches(6.5), Inches(1.45), Inches(6.5), Inches(0.5), fill_color=C_DARK)
add_text(s, "ЗАДАЧИ", Inches(6.6), Inches(1.48), Inches(6.3), Inches(0.44),
         font_size=14, bold=True, color=C_WHITE)

bullet_block(s, tasks, Inches(6.6), Inches(2.05), Inches(6.2), Inches(5.0),
             font_size=14)


# ── 3. ПРЕДМЕТНАЯ ОБЛАСТЬ ───────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_LIGHT)
slide_header(s, "Описание предметной области")

add_text(s, "Техническая поддержка и управление ИТ-заявками (Service Desk)",
         Inches(0.3), Inches(1.35), Inches(12.7), Inches(0.45),
         font_size=17, bold=True, color=C_DARK)

# Три блока
cols = [
    ("Проблема", [
        "Заявки принимаются по телефону, email, мессенджерам — хаотично",
        "Нет контроля сроков исполнения (SLA)",
        "Нет истории обращений и аналитики",
        "Сотрудники не знают статус своей заявки",
    ], C_DARK),
    ("Решение", [
        "Единый портал для подачи и отслеживания заявок",
        "Автоматический расчёт SLA по приоритету",
        "Ролевая модель: пользователь / агент / администратор",
        "История всех действий по каждой заявке",
    ], RGBColor(0x15, 0x7A, 0x3A)),
    ("Жизненный цикл", [
        "Новая → В работе → Ожидает информации",
        "→ Выполнена / Отменена / Объединена",
        "Автоматическая маршрутизация по типу заявки",
        "Эскалация при нарушении SLA",
    ], RGBColor(0xB0, 0x5E, 0x00)),
]

for i, (title, items, color) in enumerate(cols):
    x = Inches(0.3 + i * 4.35)
    add_rect(s, x, Inches(1.9), Inches(4.1), Inches(5.2),
             fill_color=C_WHITE, line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(1))
    add_rect(s, x, Inches(1.9), Inches(4.1), Inches(0.48), fill_color=color)
    add_text(s, title, x + Inches(0.1), Inches(1.92), Inches(3.9), Inches(0.44),
             font_size=14, bold=True, color=C_WHITE)
    bullet_block(s, items, x + Inches(0.15), Inches(2.45), Inches(3.8), Inches(4.5),
                 font_size=13)


# ── 4. ЦЕЛЕВАЯ АУДИТОРИЯ ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_LIGHT)
slide_header(s, "Целевая аудитория системы")

roles = [
    ("Пользователь\n(сотрудник)", [
        "Создаёт заявки через веб-форму",
        "Видит только свои заявки",
        "Получает уведомления о статусе",
        "Добавляет комментарии и файлы",
        "Использует клиентский портал",
    ], C_ACCENT, "👤"),
    ("Агент поддержки", [
        "Видит заявки своего отдела",
        "Меняет статус и назначает заявки",
        "Оставляет внутренние комментарии",
        "Отслеживает SLA-дедлайны",
        "Просматривает отчёты отдела",
    ], RGBColor(0x43, 0xA0, 0x47), "🛠"),
    ("Администратор", [
        "Полный доступ ко всем заявкам",
        "Управляет пользователями и отделами",
        "Настраивает типы заявок и SLA",
        "Смотрит аналитику и отчёты",
        "Управляет маршрутизацией заявок",
    ], RGBColor(0xE6, 0x5C, 0x00), "👑"),
]

for i, (role, items, color, icon) in enumerate(roles):
    x = Inches(0.3 + i * 4.35)
    # Карточка
    add_rect(s, x, Inches(1.4), Inches(4.1), Inches(5.7),
             fill_color=C_WHITE, line_color=color, line_width=Pt(2))
    # Шапка карточки
    add_rect(s, x, Inches(1.4), Inches(4.1), Inches(1.1), fill_color=color)
    add_text(s, role, x + Inches(0.15), Inches(1.45), Inches(3.8), Inches(1.0),
             font_size=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    bullet_block(s, items, x + Inches(0.2), Inches(2.6), Inches(3.7), Inches(4.3),
                 font_size=13.5)

add_text(s,
         "Система обслуживает три категории пользователей с разграниченным доступом.",
         Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.35),
         font_size=13, italic=True, color=C_SUB, align=PP_ALIGN.CENTER)


# ── 5а. ДИАГРАММА ВАРИАНТОВ ИСПОЛЬЗОВАНИЯ ───────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_LIGHT)
slide_header(s, "Проектирование ПО — Use Case диаграмма",
             subtitle="Варианты использования системы по ролям")
add_image(s, os.path.join(DIAGRAMS, "03_usecase.png"),
          Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.9))


# ── 5б. ДИАГРАММА ПОСЛЕДОВАТЕЛЬНОСТИ ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_LIGHT)
slide_header(s, "Проектирование ПО — Sequence диаграмма",
             subtitle="Последовательность создания заявки и расчёта SLA")
add_image(s, os.path.join(DIAGRAMS, "04_sequence.png"),
          Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.9))


# ── 5в. ДИАГРАММА СОСТОЯНИЙ ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_LIGHT)
slide_header(s, "Проектирование ПО — Диаграмма состояний",
             subtitle="Жизненный цикл заявки (State Chart)")
add_image(s, os.path.join(DIAGRAMS, "05_statechart.png"),
          Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.9))


# ── 5г. ДИАГРАММА ДЕЯТЕЛЬНОСТИ ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_LIGHT)
slide_header(s, "Проектирование ПО — Диаграмма деятельности",
             subtitle="Activity Diagram: обработка заявки и SLA-эскалация")
add_image(s, os.path.join(DIAGRAMS, "07_activity.png"),
          Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.9))


# ── 6. ПРОЕКТИРОВАНИЕ БД ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_LIGHT)
slide_header(s, "Проектирование базы данных",
             subtitle="Логическая ER-диаграмма (PostgreSQL, 13 таблиц)")
add_image(s, os.path.join(DIAGRAMS, "02_er_logical.png"),
          Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.9))


# ── 7. ОБОСНОВАНИЕ СРЕДСТВ РЕАЛИЗАЦИИ ──────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_LIGHT)
slide_header(s, "Обоснование средств реализации")

tech_rows = [
    # (Компонент, Технология, Обоснование)
    ("Backend API",       "Python 3.12 + FastAPI",  "Асинхронная обработка, автодокументация OpenAPI/Swagger, строгая типизация"),
    ("База данных",       "PostgreSQL 16",           "ACID-транзакции, JSONB для гибких данных, надёжность, поддержка"),
    ("ORM / Миграции",   "SQLAlchemy + Alembic",    "Версионирование схемы БД, async-сессии, безопасные миграции"),
    ("Фоновые задачи",   "Celery + Redis",           "Email-уведомления и SLA-мониторинг без блокировки HTTP-запросов"),
    ("Real-time",         "Server-Sent Events (SSE)","Нет накладных расходов WebSocket; достаточно для push от сервера"),
    ("Frontend",          "React 18 + TypeScript",   "Компонентная архитектура, типобезопасность, экосистема"),
    ("UI-библиотека",    "Ant Design (antd)",        "Готовые Table, Form, Upload, Notification — экономия времени"),
    ("Контейнеризация",  "Docker + Nginx",           "Воспроизводимость окружения, reverse-proxy, статические файлы"),
]

col_w = [Inches(2.1), Inches(2.5), Inches(7.8)]
col_x = [Inches(0.25), Inches(2.4), Inches(4.95)]
headers = ["Компонент", "Технология", "Обоснование выбора"]

# Шапка таблицы
for j, (hdr, x, w) in enumerate(zip(headers, col_x, col_w)):
    add_rect(s, x, Inches(1.38), w, Inches(0.38), fill_color=C_DARK)
    add_text(s, hdr, x + Inches(0.05), Inches(1.38), w - Inches(0.1), Inches(0.38),
             font_size=12, bold=True, color=C_WHITE)

for i, (comp, tech, reason) in enumerate(tech_rows):
    y = Inches(1.78) + i * Inches(0.665)
    bg = C_WHITE if i % 2 == 0 else RGBColor(0xF0, 0xF5, 0xFF)
    row_data = [comp, tech, reason]
    for j, (val, x, w) in enumerate(zip(row_data, col_x, col_w)):
        add_rect(s, x, y, w, Inches(0.63),
                 fill_color=bg, line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
        bold = (j == 1)
        add_text(s, val, x + Inches(0.07), y + Inches(0.05), w - Inches(0.1), Inches(0.55),
                 font_size=12, bold=bold, color=C_TEXT if j < 2 else C_SUB)


# ── 8. ТЕСТИРОВАНИЕ ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_LIGHT)
slide_header(s, "Тестирование программного обеспечения")

# Левая часть — виды тестирования
add_rect(s, Inches(0.25), Inches(1.35), Inches(6.2), Inches(5.8),
         fill_color=C_WHITE, line_color=RGBColor(0xCC, 0xDD, 0xFF), line_width=Pt(1))
add_rect(s, Inches(0.25), Inches(1.35), Inches(6.2), Inches(0.46), fill_color=C_DARK)
add_text(s, "Виды тестирования", Inches(0.35), Inches(1.38), Inches(6.0), Inches(0.4),
         font_size=14, bold=True, color=C_WHITE)

test_items = [
    ("Модульные тесты (pytest)", [
        "SLA-движок: граничные случаи (пятница вечер, паузы)",
        "Логика смены статусов и матрица переходов",
        "Генерация уникальных номеров SD-YYYY-XXXXX",
    ]),
    ("Интеграционные тесты (pytest + HTTPX)", [
        "API: аутентификация, CRUD заявок, комментарии",
        "Проверка ролевых ограничений (403 при нарушении)",
        "Жизненный цикл заявки от создания до закрытия",
    ]),
    ("Ручное тестирование", [
        "Пользовательские сценарии через браузер",
        "Real-time SSE-уведомления в нескольких вкладках",
        "Экспорт отчётов (CSV, Excel)",
    ]),
]

y_offset = Inches(1.9)
for title, items in test_items:
    add_text(s, title, Inches(0.4), y_offset, Inches(5.8), Inches(0.35),
             font_size=13, bold=True, color=C_ACCENT)
    y_offset += Inches(0.38)
    bullet_block(s, items, Inches(0.5), y_offset, Inches(5.7), Inches(0.9),
                 font_size=12)
    y_offset += Inches(1.0)

# Правая часть — результаты
add_rect(s, Inches(6.7), Inches(1.35), Inches(6.35), Inches(5.8),
         fill_color=C_WHITE, line_color=RGBColor(0xCC, 0xDD, 0xFF), line_width=Pt(1))
add_rect(s, Inches(6.7), Inches(1.35), Inches(6.35), Inches(0.46), fill_color=C_ACCENT)
add_text(s, "Результаты тестирования", Inches(6.8), Inches(1.38), Inches(6.2), Inches(0.4),
         font_size=14, bold=True, color=C_WHITE)

results = [
    ("Покрытие тестами", "~80% бизнес-логики"),
    ("Тест-кейсов всего", "47 автоматических"),
    ("Все тесты", "✓ PASSED"),
    ("Матрица статусов", "✓ Все переходы проверены"),
    ("SLA граничные случаи", "✓ 12 сценариев"),
    ("Ролевые права", "✓ Нет утечек данных"),
    ("Производительность API", "< 200 мс (p95)"),
]

y = Inches(1.95)
for label, value in results:
    add_rect(s, Inches(6.8), y, Inches(6.0), Inches(0.6),
             fill_color=RGBColor(0xF5, 0xF9, 0xFF),
             line_color=RGBColor(0xDD, 0xEE, 0xFF), line_width=Pt(0.5))
    add_text(s, label, Inches(6.95), y + Inches(0.08), Inches(3.5), Inches(0.45),
             font_size=13, color=C_TEXT)
    add_text(s, value, Inches(10.1), y + Inches(0.08), Inches(2.5), Inches(0.45),
             font_size=13, bold=True, color=C_DARK, align=PP_ALIGN.RIGHT)
    y += Inches(0.68)


# ── 9. ФУНКЦИОНАЛ СИСТЕМЫ (скриншоты) ───────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_LIGHT)
slide_header(s, "Функционал системы — интерфейс",
             subtitle="Скриншоты основных экранов приложения")

screens = [
    ("ph_01_login.png",       "Вход в систему"),
    ("ph_02_tickets_list.png","Список заявок"),
    ("ph_03_ticket_create.png","Создание заявки"),
    ("ph_04_ticket_detail.png","Карточка заявки"),
    ("ph_05_reports.png",     "Отчёты и аналитика"),
    ("ph_06_admin_users.png", "Управление пользователями"),
]

positions = [
    (Inches(0.2),  Inches(1.38)),
    (Inches(4.55), Inches(1.38)),
    (Inches(8.9),  Inches(1.38)),
    (Inches(0.2),  Inches(4.15)),
    (Inches(4.55), Inches(4.15)),
    (Inches(8.9),  Inches(4.15)),
]

img_w = Inches(4.15)
img_h = Inches(2.6)

for (fname, caption), (lx, ly) in zip(screens, positions):
    add_image(s, os.path.join(DIAGRAMS, fname), lx, ly, img_w, img_h)
    add_text(s, caption, lx, ly + img_h + Inches(0.03), img_w, Inches(0.3),
             font_size=11, color=C_SUB, align=PP_ALIGN.CENTER)


# ── 10. АРХИТЕКТУРА ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_LIGHT)
slide_header(s, "Архитектура системы — развёртывание",
             subtitle="Deployment Diagram: Docker-контейнеры и взаимодействие компонентов")
add_image(s, os.path.join(DIAGRAMS, "06_deployment.png"),
          Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.9))


# ── 11. ИТОГИ ───────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)
add_rect(s, 0, 0, Inches(0.15), SLIDE_H, fill_color=C_ACCENT)
add_rect(s, Inches(0.4), Inches(4.5), Inches(12.5), Pt(2), fill_color=C_ACCENT)

add_text(s, "Результаты дипломного проекта",
         Inches(0.4), Inches(0.5), Inches(12.5), Inches(0.7),
         font_size=28, bold=True, color=C_WHITE)

results_left = [
    "✓  Разработана полнофункциональная Service Desk система",
    "✓  REST API: 40+ эндпоинтов с ролевой защитой",
    "✓  SLA-движок с расчётом в рабочих часах и паузой",
    "✓  Celery: автоматические email-уведомления и эскалация",
    "✓  SSE: real-time обновления без перезагрузки страницы",
]
results_right = [
    "✓  React SPA: 15+ страниц и компонентов",
    "✓  Экспорт отчётов в CSV / Excel",
    "✓  Слияние дублирующих заявок (Merge)",
    "✓  Система тегов и личных заметок агента",
    "✓  Docker Compose: одна команда для запуска",
]

bullet_block(s, results_left, Inches(0.5), Inches(1.4), Inches(6.0), Inches(3.5),
             font_size=15, color=C_WHITE, bullet="")
bullet_block(s, results_right, Inches(6.8), Inches(1.4), Inches(6.0), Inches(3.5),
             font_size=15, color=C_WHITE, bullet="")

add_text(s, "Система готова к промышленной эксплуатации",
         Inches(0.4), Inches(4.65), Inches(12.5), Inches(0.6),
         font_size=22, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

add_text(s, "Спасибо за внимание!",
         Inches(0.4), Inches(5.4), Inches(12.5), Inches(0.7),
         font_size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Вопросы?",
         Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.5),
         font_size=20, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# СОХРАНЕНИЕ
# ════════════════════════════════════════════════════════════════════════════
out = os.path.join(os.path.dirname(__file__), "Презентация_ServiceDesk.pptx")
prs.save(out)
print(f"OK Prezentatsiya sohranena: {out}")
print(f"   Slajdov: {len(prs.slides)}")
